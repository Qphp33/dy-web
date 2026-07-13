from pathlib import Path
import sys

sys.path.insert(0, str(Path(__file__).resolve().parent / ".ppt_deps"))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path("苏教版三下_统计表和条形统计图_互动练习.pptx")


BLUE = RGBColor(43, 117, 221)
LIGHT_BLUE = RGBColor(232, 242, 255)
DEEP_BLUE = RGBColor(30, 76, 153)
TEXT = RGBColor(34, 51, 84)
GRAY = RGBColor(102, 119, 153)
GREEN = RGBColor(82, 196, 26)
RED = RGBColor(245, 34, 45)
ORANGE = RGBColor(250, 140, 22)
WHITE = RGBColor(255, 255, 255)


def add_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=""):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.5), Inches(0.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Microsoft YaHei"
    r.font.bold = True
    r.font.size = Pt(24)
    r.font.color.rgb = DEEP_BLUE
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.65), Inches(1.0), Inches(8.5), Inches(0.4))
        tf2 = sub.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Microsoft YaHei"
        r2.font.size = Pt(11)
        r2.font.color.rgb = GRAY


def add_round_box(slide, left, top, width, height, text="", fill_color=LIGHT_BLUE, line_color=BLUE, font_size=16, bold=False):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1.2)
    if text:
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(font_size)
        r.font.bold = bold
        r.font.color.rgb = TEXT
    return shape


def add_bullets(slide, items, left=0.9, top=1.6, width=8.2, height=4.7, font_size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)
        p.bullet = True


def add_footer(slide, text="苏教版三年级下册数学"):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.1), Inches(10), Inches(0.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(10)
    r.font.color.rgb = WHITE


def add_table_like(slide, headers, rows, left=0.9, top=1.8, col_widths=None):
    if col_widths is None:
        col_widths = [1.8] * len(headers)
    x = Inches(left)
    y = Inches(top)
    row_h = Inches(0.55)
    for i, head in enumerate(headers):
        w = Inches(col_widths[i])
        add_round_box(slide, x, y, w, row_h, head, BLUE, BLUE, 14, True)
        slide.shapes[-1].text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
        x += w
    for r_idx, row in enumerate(rows):
        x = Inches(left)
        y2 = y + row_h * (r_idx + 1)
        for c_idx, cell in enumerate(row):
            w = Inches(col_widths[c_idx])
            add_round_box(slide, x, y2, w, row_h, str(cell), WHITE, RGBColor(190, 214, 245), 14, False)
            x += w


def add_chart_axes(slide, left, top, width, height, y_max=10, y_step=2, x_labels=None):
    x_labels = x_labels or []
    l = Inches(left)
    t = Inches(top)
    w = Inches(width)
    h = Inches(height)
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, l, t + h - Pt(2), w, Pt(2))
    slide.shapes[-1].fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = DEEP_BLUE
    slide.shapes[-1].line.fill.background()
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, l, t, Pt(2), h)
    slide.shapes[-1].fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = DEEP_BLUE
    slide.shapes[-1].line.fill.background()
    grid_count = y_max // y_step
    for i in range(grid_count + 1):
        yy = t + h - (h / grid_count) * i if grid_count else t
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, l, yy, w, Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(221, 233, 249)
        line.line.fill.background()
        box = slide.shapes.add_textbox(l - Inches(0.35), yy - Inches(0.1), Inches(0.3), Inches(0.2))
        p = box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = str(i * y_step)
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(10)
        r.font.color.rgb = GRAY
    if x_labels:
        step = w / len(x_labels)
        for idx, lab in enumerate(x_labels):
            box = slide.shapes.add_textbox(l + step * idx + step * 0.15, t + h + Inches(0.05), step * 0.7, Inches(0.25))
            p = box.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = lab
            r.font.name = "Microsoft YaHei"
            r.font.size = Pt(11)
            r.font.color.rgb = TEXT


def add_bars(slide, left, top, width, height, values, colors):
    l = Inches(left)
    t = Inches(top)
    w = Inches(width)
    h = Inches(height)
    count = len(values)
    gap = w * 0.12 / count
    bar_w = (w - gap * (count + 1)) / count
    max_v = max(values)
    for idx, val in enumerate(values):
        bh = h * val / max_v
        x = l + gap * (idx + 1) + bar_w * idx
        y = t + h - bh
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, bar_w, bh)
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[idx]
        shape.line.color.rgb = colors[idx]
        box = slide.shapes.add_textbox(x, y - Inches(0.22), bar_w, Inches(0.2))
        p = box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(val)
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = colors[idx]


def add_link_button(slide, text, left, top, width, height, target_slide, fill_color, font_color=WHITE):
    shape = add_round_box(slide, left, top, width, height, text, fill_color, fill_color, 16, True)
    shape.text_frame.paragraphs[0].runs[0].font.color.rgb = font_color
    shape.click_action.target_slide = target_slide
    return shape


prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
layout = prs.slide_layouts[6]

# 1 cover
slide = prs.slides.add_slide(layout)
add_bg(slide, WHITE)
cover = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.55), Inches(8.9), Inches(5.9))
cover.fill.solid()
cover.fill.fore_color.rgb = LIGHT_BLUE
cover.line.color.rgb = RGBColor(190, 214, 245)
banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(1.0), Inches(3.4), Inches(0.6))
banner.fill.solid()
banner.fill.fore_color.rgb = BLUE
banner.line.fill.background()
tf = banner.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "苏教版三年级下册"
r.font.name = "Microsoft YaHei"
r.font.size = Pt(18)
r.font.bold = True
r.font.color.rgb = WHITE
t1 = slide.shapes.add_textbox(Inches(0.95), Inches(1.95), Inches(7.5), Inches(1.2))
p = t1.text_frame.paragraphs[0]
r = p.add_run()
r.text = "统计表和条形统计图"
r.font.name = "Microsoft YaHei"
r.font.bold = True
r.font.size = Pt(28)
r.font.color.rgb = DEEP_BLUE
t2 = slide.shapes.add_textbox(Inches(0.95), Inches(3.0), Inches(7.8), Inches(1.2))
for idx, txt in enumerate(["认识统计表", "会看条形统计图", "互动练习: 拖拽数据块并即时反馈"]):
    p = t2.text_frame.paragraphs[0] if idx == 0 else t2.text_frame.add_paragraph()
    p.text = txt
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT
    p.space_after = Pt(10)
    p.bullet = True
add_footer(slide)

# 2 objectives
slide = prs.slides.add_slide(layout)
add_bg(slide)
add_title(slide, "学习目标", "看懂统计表，学会根据数据绘制和分析条形统计图")
items = [
    "能从生活情境中收集和整理数据。",
    "会用统计表表示数据，并说出每类数量。",
    "会根据统计表完成简单条形统计图。",
    "能根据统计图比较多少，并提出简单数学问题。",
]
add_bullets(slide, items)
add_footer(slide)

# 3 scenario
slide = prs.slides.add_slide(layout)
add_bg(slide)
add_title(slide, "创设情境", "三年级同学最喜欢的水果调查")
add_round_box(slide, Inches(0.8), Inches(1.6), Inches(8.4), Inches(1.0), "老师调查了全班同学最喜欢的水果，结果如下：苹果 8 人，香蕉 6 人，橙子 4 人，葡萄 2 人。", WHITE, RGBColor(190, 214, 245), 18, False)
fruit_colors = [RGBColor(255, 77, 79), RGBColor(250, 173, 20), RGBColor(255, 160, 0), RGBColor(114, 46, 209)]
labels = ["苹果", "香蕉", "橙子", "葡萄"]
for idx, lab in enumerate(labels):
    add_round_box(slide, Inches(1.0 + idx * 2.1), Inches(3.0), Inches(1.6), Inches(1.05), f"{lab}\n{[8,6,4,2][idx]}人", fruit_colors[idx], fruit_colors[idx], 16, True).text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
add_bullets(slide, ["想一想：怎样把这些数据整理得更清楚？", "我们可以先制成统计表，再画条形统计图。"], left=1.0, top=4.6, width=8.0, height=1.5, font_size=18)
add_footer(slide)

# 4 table
slide = prs.slides.add_slide(layout)
add_bg(slide)
add_title(slide, "认识统计表", "把调查得到的数据分类整理")
add_table_like(slide, ["水果", "人数"], [["苹果", 8], ["香蕉", 6], ["橙子", 4], ["葡萄", 2]], left=1.5, top=1.8, col_widths=[3.2, 2.0])
add_bullets(slide, ["统计表能清楚地看出每一类有多少。", "看统计表时，要注意：项目名称、数量、单位。", "从表中看出：喜欢苹果的人最多，喜欢葡萄的人最少。"], left=0.95, top=4.8, width=8.2, height=1.8, font_size=18)
add_footer(slide)

# 5 chart
slide = prs.slides.add_slide(layout)
add_bg(slide)
add_title(slide, "认识条形统计图", "用条形的高低表示数量的多少")
add_chart_axes(slide, 1.1, 1.7, 5.6, 4.2, y_max=10, y_step=2, x_labels=labels)
add_bars(slide, 1.1, 1.7, 5.6, 4.2, [8, 6, 4, 2], fruit_colors)
add_bullets(slide, ["横轴写项目，纵轴写数量。", "条形越高，表示数量越多。", "画图时每一格表示的数量要统一。"], left=7.0, top=2.0, width=2.3, height=3.6, font_size=18)
add_footer(slide)

# 6 method
slide = prs.slides.add_slide(layout)
add_bg(slide)
add_title(slide, "画图步骤", "根据统计表画条形统计图")
steps = [
    "第一步：先看清统计表中的项目和数量。",
    "第二步：在横轴上写出项目名称。",
    "第三步：在纵轴上按统一标准标出数量。",
    "第四步：按数量画出对应高度的条形。",
    "第五步：检查每个条形高度是否正确。",
]
for i, s in enumerate(steps):
    card = add_round_box(slide, Inches(0.95), Inches(1.5 + i * 0.95), Inches(8.1), Inches(0.68), s, WHITE, RGBColor(190, 214, 245), 18, False)
    num = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.55), Inches(1.54 + i * 0.95), Inches(0.45), Inches(0.45))
    num.fill.solid()
    num.fill.fore_color.rgb = BLUE
    num.line.fill.background()
    p = num.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(i + 1)
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = WHITE
add_footer(slide)

# 7 drag activity
slide = prs.slides.add_slide(layout)
add_bg(slide)
add_title(slide, "互动练习一", "拖拽数据块，完成统计图")
tip = add_round_box(slide, Inches(0.8), Inches(1.15), Inches(8.7), Inches(0.6), "玩法：先看左侧统计表，再把右侧下方的彩色数据块拖到对应位置，拼出完整条形统计图。", WHITE, RGBColor(190, 214, 245), 15, False)
tip.text_frame.paragraphs[0].runs[0].font.color.rgb = GRAY
add_table_like(slide, ["项目", "本数"], [["故事书", 5], ["科技书", 7], ["漫画书", 3], ["童话书", 6]], left=0.8, top=1.95, col_widths=[1.6, 1.2])
add_chart_axes(slide, 3.6, 2.0, 4.9, 3.6, y_max=8, y_step=1, x_labels=["故事书", "科技书", "漫画书", "童话书"])
for idx, txt in enumerate(["5", "7", "3", "6"]):
    shape = add_round_box(slide, Inches(3.9 + idx * 1.18), Inches(6.05), Inches(0.85), Inches(0.42), txt, [RGBColor(105, 192, 255), RGBColor(54, 207, 201), RGBColor(255, 169, 64), RGBColor(177, 107, 255)][idx], [RGBColor(105, 192, 255), RGBColor(54, 207, 201), RGBColor(255, 169, 64), RGBColor(177, 107, 255)][idx], 18, True)
    shape.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
note = slide.shapes.add_textbox(Inches(0.8), Inches(6.1), Inches(2.7), Inches(0.55))
p = note.text_frame.paragraphs[0]
r = p.add_run()
r.text = "提示：本页为可编辑拖拽练习。"
r.font.name = "Microsoft YaHei"
r.font.size = Pt(14)
r.font.bold = True
r.font.color.rgb = ORANGE
add_footer(slide)

# 8 choose correct
slide8 = prs.slides.add_slide(layout)
add_bg(slide8)
add_title(slide8, "互动练习二", "看统计表，选出正确的条形统计图")
add_table_like(slide8, ["动物", "只数"], [["小兔", 4], ["小猫", 7], ["小狗", 5]], left=0.8, top=1.7, col_widths=[1.5, 1.2])
add_round_box(slide8, Inches(3.0), Inches(1.55), Inches(6.0), Inches(0.5), "请点击你认为正确的一幅统计图", WHITE, RGBColor(190, 214, 245), 16, True)
for idx, option in enumerate(["A", "B", "C"]):
    left = 3.2 + idx * 1.9
    add_round_box(slide8, Inches(left), Inches(2.3), Inches(1.45), Inches(2.45), option, RGBColor(248, 251, 255), RGBColor(190, 214, 245), 22, True)
    add_chart_axes(slide8, left + 0.1, 2.7, 1.25, 1.7, y_max=8, y_step=2, x_labels=["兔", "猫", "狗"])
    if option == "A":
        vals = [4, 7, 5]
    elif option == "B":
        vals = [4, 5, 7]
    else:
        vals = [7, 4, 5]
    add_bars(slide8, left + 0.1, 2.7, 1.25, 1.7, vals, [BLUE, GREEN, ORANGE])

# 9 right feedback
slide9 = prs.slides.add_slide(layout)
add_bg(slide9, RGBColor(246, 255, 237))
add_title(slide9, "回答正确", "你真棒，选对了")
ok = slide9.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(3.9), Inches(2.0), Inches(1.2), Inches(1.2))
ok.fill.solid()
ok.fill.fore_color.rgb = GREEN
ok.line.fill.background()
p = ok.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "√"
r.font.name = "Microsoft YaHei"
r.font.size = Pt(34)
r.font.bold = True
r.font.color.rgb = WHITE
add_bullets(slide9, ["正确图应表示：小兔 4 只，小猫 7 只，小狗 5 只。", "条形统计图要和统计表一一对应。"], left=2.0, top=3.6, width=6.1, height=1.5, font_size=20)
add_footer(slide9)

# 10 wrong feedback
slide10 = prs.slides.add_slide(layout)
add_bg(slide10, RGBColor(255, 241, 240))
add_title(slide10, "再想一想", "别着急，我们一起检查")
no = slide10.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(3.9), Inches(2.0), Inches(1.2), Inches(1.2))
no.fill.solid()
no.fill.fore_color.rgb = RED
no.line.fill.background()
p = no.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "×"
r.font.name = "Microsoft YaHei"
r.font.size = Pt(34)
r.font.bold = True
r.font.color.rgb = WHITE
add_bullets(slide10, ["回到统计表再看一看：哪一类最多？哪一类最少？", "比较每个条形的高低是否和表中的数量一致。"], left=1.6, top=3.6, width=6.8, height=1.5, font_size=20)
add_footer(slide10)

# 11 summary
slide = prs.slides.add_slide(layout)
add_bg(slide)
add_title(slide, "课堂小结", "今天你学会了什么？")
add_bullets(slide, ["会看统计表，知道每一类的数据。", "会根据统计表画条形统计图。", "会根据统计图比较多少，并解决简单问题。", "记住：统计让数据更清楚，图表让信息更直观。"])
add_footer(slide)

# links after target slides exist
for shape in slide8.shapes:
    if not shape.has_text_frame:
        continue
    text = shape.text_frame.text.strip()
    if text == "A":
        shape.click_action.target_slide = slide9
    elif text in {"B", "C"}:
        shape.click_action.target_slide = slide10

back1 = add_link_button(slide9, "返回练习", Inches(3.8), Inches(5.5), Inches(1.6), Inches(0.55), slide8, BLUE)
back2 = add_link_button(slide10, "返回练习", Inches(3.8), Inches(5.5), Inches(1.6), Inches(0.55), slide8, BLUE)

prs.save(OUT)
print(str(OUT.resolve()))
